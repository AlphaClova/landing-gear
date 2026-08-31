from __future__ import annotations

import argparse
import hashlib
import json
import re
from pathlib import Path
from typing import Any

import pdfplumber
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


SUPPORTED_TYPES = {"pdf", "docx", "xlsx", "pptx"}
METADATA_MAP: dict[str, dict[str, Any]] = {
    "doc10": {
        "topic": "pension_system",
        "account_type": "DB_DC",
        "effective_from": None,
        "valid_to": None,
    },
    "doc51": {
        "topic": "withdrawal_tax",
        "account_type": "IRP",
        # The file has authoring timestamps, but no stated effective date.
        "effective_from": None,
        "valid_to": None,
    },
    "doc26": {
        "topic": "withdrawal_tax",
        "account_type": "IRP",
        "effective_from": None,
        "valid_to": None,
    },
    "doc41": {
        "topic": "withdrawal_tax",
        "account_type": "IRP",
        "effective_from": None,
        "valid_to": None,
    },
    "doc55": {
        "topic": "withdrawal_tax",
        "account_type": "IRP",
        "effective_from": None,
        "valid_to": None,
    },
    **{
        document_id: {
            "topic": "product",
            "account_type": "IRP",
            "effective_from": None,
            "valid_to": None,
        }
        for document_id in (
            "r2_kr5153420063",
            "r2_kr5153420079",
            "r2_kr5153420105",
            "r2_kr5153450658",
        )
    },
}
PROVIDED_SOURCE_PRIORITY = 0
DEFAULT_TARGET_DOCUMENTS = tuple(METADATA_MAP)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Create document inventory and processed chunks from raw PDF/DOCX/XLSX files."
    )
    parser.add_argument("--raw-dir", default="app/data/raw", help="Raw document directory")
    parser.add_argument(
        "--inventory-output",
        default="app/data/processed/document_inventory.json",
        help="Document inventory output path",
    )
    parser.add_argument(
        "--chunks-output",
        default="app/data/processed/chunks.jsonl",
        help="Chunk output JSONL path",
    )
    parser.add_argument(
        "--target-documents",
        default=",".join(DEFAULT_TARGET_DOCUMENTS),
        help="Comma-separated document ids to parse into chunk output",
    )
    return parser.parse_args()


def chunk_text(text: str, chunk_size: int = 900) -> list[str]:
    lines = [line.strip() for line in text.splitlines() if line.strip()]
    chunks: list[str] = []
    current = ""

    for line in lines:
        candidate = f"{current} {line}".strip()
        if len(candidate) > chunk_size and current:
            chunks.append(current)
            current = line
        else:
            current = candidate

    if current:
        chunks.append(current)
    return chunks


def detect_document_id(file_path: Path) -> str:
    match = re.search(r"(doc\d+)", file_path.stem.lower())
    if match:
        return match.group(1)
    return file_path.stem.lower()


def resolve_metadata(document_id: str) -> dict[str, Any]:
    metadata = {
        # Do not infer metadata for inventory-only documents from their filename.
        "topic": None,
        "account_type": None,
        "effective_from": None,
        "valid_to": None,
    }
    metadata.update(METADATA_MAP.get(document_id, {}))
    return metadata


def extract_docx_rendered_page_count(document: Document) -> int | None:
    """Return a page count only when the DOCX stores rendered page boundaries."""
    rendered_breaks = len(document.element.body.xpath(".//w:lastRenderedPageBreak"))
    explicit_breaks = len(document.element.body.xpath('.//w:br[@w:type="page"]'))
    breaks = rendered_breaks or explicit_breaks
    return breaks + 1 if breaks else None


def _first_non_empty(lines: list[str]) -> str:
    for line in lines:
        line_stripped = line.strip()
        if line_stripped:
            return line_stripped
    return ""


def _normalize_table_rows(rows: list[list[str | None]]) -> str:
    normalized: list[str] = []
    for row in rows:
        cells = [(cell or "").strip() for cell in row]
        normalized.append(" | ".join(cells))
    return "\n".join(normalized).strip()


def _chunk_id(document_id: str, page: int | None, content_type: str, index: int, content: str) -> str:
    digest = hashlib.sha1(content.encode("utf-8")).hexdigest()[:10]
    page_token = "na" if page is None else f"{page:03d}"
    return f"{document_id}-p{page_token}-{content_type[0]}{index:03d}-{digest}"


def _parse_quality(content: str, content_type: str) -> str:
    if not content.strip():
        return "low"
    if content_type == "table" and content.count("|") >= 2:
        return "high"
    if len(content) >= 80:
        return "high"
    return "medium"


def _build_chunk(
    document_id: str,
    title: str,
    page: int | None,
    section: str,
    topic: str | None,
    account_type: str | None,
    effective_from: str | None,
    valid_to: str | None,
    source_priority: int,
    content_type: str,
    content: str,
    index: int,
) -> dict[str, Any]:
    return {
        "chunk_id": _chunk_id(document_id, page, content_type, index, content),
        "document_id": document_id,
        "title": title,
        "page": page,
        "section": section,
        "text": content,
        "content": content,
        "content_type": content_type,
        "parse_quality": _parse_quality(content, content_type),
        "effective_from": effective_from,
        "valid_to": valid_to,
        "topics": [topic] if topic is not None else [],
        "topic": topic,
        "account_types": [account_type] if account_type is not None else [],
        "account_type": account_type,
        "source_type": "provided",
        "source_priority": source_priority,
    }


def parse_pdf(file_path: Path, document_id: str, metadata: dict[str, Any]) -> tuple[str, int, list[dict[str, Any]]]:
    chunks: list[dict[str, Any]] = []
    with pdfplumber.open(file_path) as pdf:
        page_count = len(pdf.pages)
        first_page_text = ""
        for page_index, page in enumerate(pdf.pages, start=1):
            page_text = (page.extract_text() or "").strip()
            if not first_page_text and page_text:
                first_page_text = page_text

            section = _first_non_empty(page_text.splitlines()) or f"page-{page_index}"
            text_chunks = chunk_text(page_text)
            for chunk_index, piece in enumerate(text_chunks, start=1):
                chunks.append(
                    _build_chunk(
                        document_id=document_id,
                        title=file_path.stem,
                        page=page_index,
                        section=section,
                        topic=metadata["topic"],
                        account_type=metadata["account_type"],
                        effective_from=metadata["effective_from"],
                        valid_to=metadata["valid_to"],
                        source_priority=PROVIDED_SOURCE_PRIORITY,
                        content_type="text",
                        content=piece,
                        index=chunk_index,
                    )
                )

            tables = page.extract_tables() or []
            for table_index, table_rows in enumerate(tables, start=1):
                table_content = _normalize_table_rows(table_rows)
                if not table_content:
                    continue
                chunks.append(
                    _build_chunk(
                        document_id=document_id,
                        title=file_path.stem,
                        page=page_index,
                        section=f"{section}-table-{table_index}",
                        topic=metadata["topic"],
                        account_type=metadata["account_type"],
                        effective_from=metadata["effective_from"],
                        valid_to=metadata["valid_to"],
                        source_priority=PROVIDED_SOURCE_PRIORITY,
                        content_type="table",
                        content=table_content,
                        index=table_index,
                    )
                )

        title = _first_non_empty(first_page_text.splitlines()) or file_path.stem
    return title, page_count, chunks


def parse_docx(file_path: Path, document_id: str, metadata: dict[str, Any]) -> tuple[str, int | None, list[dict[str, Any]]]:
    document = Document(file_path)
    title = _first_non_empty([paragraph.text for paragraph in document.paragraphs]) or file_path.stem
    page_count = extract_docx_rendered_page_count(document)

    chunks: list[dict[str, Any]] = []
    current_page = 1 if page_count is not None else None
    current_section = title
    text_index = 0
    table_index = 0

    # Iterate body blocks in source order so section and rendered-page provenance
    # remain traceable. A page is emitted only when the DOCX itself stores breaks.
    for child in document.element.body.iterchildren():
        if child.tag.endswith("}p"):
            paragraph = next((p for p in document.paragraphs if p._p is child), None)
            if paragraph is None:
                continue
            break_count = len(child.xpath(".//w:lastRenderedPageBreak"))
            if not break_count:
                break_count = len(child.xpath('.//w:br[@w:type="page"]'))
            if current_page is not None:
                current_page += break_count
            content = paragraph.text.strip()
            if not content:
                continue
            if paragraph.style.name == "List Paragraph":
                current_section = content
            text_index += 1
            chunks.append(
                _build_chunk(
                    document_id=document_id,
                    title=title,
                    page=current_page,
                    section=current_section,
                    topic=metadata["topic"],
                    account_type=metadata["account_type"],
                    effective_from=metadata["effective_from"],
                    valid_to=metadata["valid_to"],
                    source_priority=PROVIDED_SOURCE_PRIORITY,
                    content_type="text",
                    content=content,
                    index=text_index,
                )
            )
        elif child.tag.endswith("}tbl"):
            table = next((t for t in document.tables if t._tbl is child), None)
            if table is None:
                continue
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            table_content = _normalize_table_rows(rows)
            if table_content:
                table_index += 1
                chunks.append(
                    _build_chunk(
                        document_id=document_id,
                        title=title,
                        page=current_page,
                        section=current_section,
                        topic=metadata["topic"],
                        account_type=metadata["account_type"],
                        effective_from=metadata["effective_from"],
                        valid_to=metadata["valid_to"],
                        source_priority=PROVIDED_SOURCE_PRIORITY,
                        content_type="table",
                        content=table_content,
                        index=table_index,
                    )
                )

    return title, page_count, chunks


def parse_xlsx(file_path: Path, document_id: str, metadata: dict[str, Any]) -> tuple[str, list[dict[str, Any]], list[dict[str, Any]]]:
    workbook = load_workbook(file_path, read_only=True, data_only=True)
    sheet_info: list[dict[str, Any]] = []
    chunks: list[dict[str, Any]] = []

    for sheet_index, sheet_name in enumerate(workbook.sheetnames, start=1):
        sheet = workbook[sheet_name]
        sheet_info.append(
            {
                "name": sheet_name,
                "max_row": sheet.max_row,
                "max_column": sheet.max_column,
            }
        )

        rows: list[list[str | None]] = []
        for row in sheet.iter_rows(values_only=True):
            row_values: list[str | None] = []
            for cell in row:
                row_values.append(None if cell is None else str(cell))
            if any(value not in (None, "") for value in row_values):
                rows.append(row_values)

        if not rows:
            continue

        table_content = _normalize_table_rows(rows)
        chunks.append(
            _build_chunk(
                document_id=document_id,
                title=file_path.stem,
                page=sheet_index,
                section=sheet_name,
                topic=metadata["topic"],
                account_type=metadata["account_type"],
                effective_from=metadata["effective_from"],
                valid_to=metadata["valid_to"],
                source_priority=PROVIDED_SOURCE_PRIORITY,
                content_type="table",
                content=table_content,
                index=1,
            )
        )

    workbook.close()
    title = file_path.stem
    return title, sheet_info, chunks


def _slide_section(slide: Any) -> str | None:
    """Use only source text as a section; never synthesize a slide heading."""
    if slide.shapes.title is not None:
        title = slide.shapes.title.text.replace("\v", "\n").strip()
        if title:
            return title
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.replace("\v", "\n").strip()
            if text:
                return text
    return None


def parse_pptx(
    file_path: Path, document_id: str, metadata: dict[str, Any]
) -> tuple[str, int, list[dict[str, Any]]]:
    presentation = Presentation(file_path)
    chunks: list[dict[str, Any]] = []
    document_title = file_path.stem

    for slide_number, slide in enumerate(presentation.slides, start=1):
        section = _slide_section(slide)
        if slide_number == 1 and section:
            document_title = section

        text_parts: list[str] = []
        has_image = False
        table_index = 0
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                has_image = True
            if shape.has_table:
                rows = [
                    [cell.text.replace("\v", "\n") for cell in row.cells]
                    for row in shape.table.rows
                ]
                table_content = _normalize_table_rows(rows)
                if not table_content:
                    continue
                table_index += 1
                table_chunk = _build_chunk(
                    document_id=document_id,
                    title=document_title,
                    # For PPTX, page is the stable 1-based slide number.
                    page=slide_number,
                    section=section or file_path.stem,
                    topic=metadata["topic"],
                    account_type=metadata["account_type"],
                    effective_from=metadata["effective_from"],
                    valid_to=metadata["valid_to"],
                    source_priority=PROVIDED_SOURCE_PRIORITY,
                    content_type="table",
                    content=table_content,
                    index=table_index,
                )
                chunks.append(table_chunk)
            elif shape.has_text_frame:
                text = shape.text.replace("\v", "\n").strip()
                if text:
                    text_parts.append(text)

        content = "\n".join(text_parts).strip()
        if content:
            text_chunk = _build_chunk(
                document_id=document_id,
                title=document_title,
                page=slide_number,
                section=section or file_path.stem,
                topic=metadata["topic"],
                account_type=metadata["account_type"],
                effective_from=metadata["effective_from"],
                valid_to=metadata["valid_to"],
                source_priority=PROVIDED_SOURCE_PRIORITY,
                content_type="text",
                content=content,
                index=1,
            )
            if has_image and len(content) < 80:
                text_chunk["parse_quality"] = "low"
            chunks.append(text_chunk)

    return document_title, len(presentation.slides), chunks


def parse_document(file_path: Path) -> tuple[dict[str, Any], list[dict[str, Any]]]:
    document_id = detect_document_id(file_path)
    file_type = file_path.suffix.lower().lstrip(".")
    metadata = resolve_metadata(document_id)

    inventory_entry: dict[str, Any] = {
        "document_id": document_id,
        "filename": file_path.name,
        "raw_path": "",
        "file_type": file_type,
        "page_count": None,
        "sheet_info": None,
        "title": file_path.stem,
        "topic": metadata["topic"],
        "account_type": metadata["account_type"],
        "source_priority": PROVIDED_SOURCE_PRIORITY,
        "effective_from": metadata["effective_from"],
        "valid_to": metadata["valid_to"],
        "parse_status": "inventoried",
    }

    if file_type not in SUPPORTED_TYPES:
        inventory_entry["parse_status"] = "unsupported_file_type"
        return inventory_entry, []

    if file_type == "pdf":
        title, page_count, chunks = parse_pdf(file_path, document_id, metadata)
        inventory_entry["title"] = title
        inventory_entry["page_count"] = page_count
        return inventory_entry, chunks

    if file_type == "docx":
        title, page_count, chunks = parse_docx(file_path, document_id, metadata)
        inventory_entry["title"] = title
        inventory_entry["page_count"] = page_count
        return inventory_entry, chunks

    if file_type == "pptx":
        title, slide_count, chunks = parse_pptx(file_path, document_id, metadata)
        inventory_entry["title"] = title
        inventory_entry["page_count"] = slide_count
        return inventory_entry, chunks

    title, sheet_info, chunks = parse_xlsx(file_path, document_id, metadata)
    inventory_entry["title"] = title
    inventory_entry["sheet_info"] = sheet_info
    return inventory_entry, chunks


def build_inventory(raw_dir: Path) -> list[dict[str, Any]]:
    entries: list[dict[str, Any]] = []
    for file_path in sorted(path for path in raw_dir.rglob("*") if path.is_file()):
        doc_id = detect_document_id(file_path)
        meta = resolve_metadata(doc_id)
        file_type = file_path.suffix.lower().lstrip(".")
        relative_path = file_path.relative_to(raw_dir).as_posix()

        entry: dict[str, Any] = {
            "document_id": doc_id,
            "filename": file_path.name,
            "raw_path": relative_path,
            "file_type": file_type,
            "page_count": None,
            "sheet_info": None,
            "title": file_path.stem,
            "topic": meta["topic"],
            "account_type": meta["account_type"],
            "source_priority": PROVIDED_SOURCE_PRIORITY,
            "effective_from": meta["effective_from"],
            "valid_to": meta["valid_to"],
            "parse_status": "inventoried",
        }

        try:
            if file_type == "pdf":
                with pdfplumber.open(file_path) as pdf:
                    entry["page_count"] = len(pdf.pages)
                    if pdf.pages:
                        first_text = (pdf.pages[0].extract_text() or "").splitlines()
                        first_line = _first_non_empty(first_text)
                        if first_line:
                            entry["title"] = first_line
            elif file_type == "docx":
                document = Document(file_path)
                entry["page_count"] = extract_docx_rendered_page_count(document)
                first_line = _first_non_empty([paragraph.text for paragraph in document.paragraphs])
                if first_line:
                    entry["title"] = first_line
            elif file_type == "xlsx":
                workbook = load_workbook(file_path, read_only=True, data_only=True)
                entry["sheet_info"] = [
                    {
                        "name": sheet_name,
                        "max_row": workbook[sheet_name].max_row,
                        "max_column": workbook[sheet_name].max_column,
                    }
                    for sheet_name in workbook.sheetnames
                ]
                workbook.close()
            elif file_type == "pptx":
                presentation = Presentation(file_path)
                entry["page_count"] = len(presentation.slides)
                if presentation.slides:
                    first_section = _slide_section(presentation.slides[0])
                    if first_section:
                        entry["title"] = first_section
            else:
                entry["parse_status"] = "unsupported_file_type"
        except Exception as exc:  # pragma: no cover - real file variance
            entry["parse_status"] = "inventory_failed"
            entry["parse_error"] = str(exc)

        entries.append(entry)
    return entries


def parse_targets(
    raw_dir: Path,
    inventory: list[dict[str, Any]],
    target_documents: set[str],
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], list[dict[str, Any]]]:
    chunks: list[dict[str, Any]] = []
    parsed_documents: list[dict[str, Any]] = []
    failed_documents: list[dict[str, Any]] = []

    indexed = {entry["document_id"]: entry for entry in inventory}
    available_doc_ids = set(indexed)

    for target_id in sorted(target_documents):
        if target_id not in available_doc_ids:
            failed = {
                "document_id": target_id,
                "filename": None,
                "parse_status": "missing_source",
                "parse_error": "raw source file not found",
            }
            failed_documents.append(failed)
            continue

        source_entry = indexed[target_id]
        file_path = raw_dir / source_entry["raw_path"]
        try:
            parsed_entry, parsed_chunks = parse_document(file_path)
            source_entry["title"] = parsed_entry["title"]
            source_entry["page_count"] = parsed_entry["page_count"]
            source_entry["sheet_info"] = parsed_entry["sheet_info"]
            source_entry["parse_status"] = "parsed"

            chunks.extend(parsed_chunks)
            parsed_documents.append(
                {
                    "document_id": target_id,
                    "filename": source_entry["filename"],
                    "chunk_count": len(parsed_chunks),
                }
            )
        except Exception as exc:  # pragma: no cover - real file variance
            source_entry["parse_status"] = "parse_failed"
            source_entry["parse_error"] = str(exc)
            failed_documents.append(
                {
                    "document_id": target_id,
                    "filename": source_entry["filename"],
                    "parse_status": "parse_failed",
                    "parse_error": str(exc),
                }
            )

    return inventory, chunks, failed_documents


def write_inventory(path: Path, inventory: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(inventory, ensure_ascii=True, indent=2), encoding="utf-8")


def write_chunks(path: Path, chunks: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as fp:
        for chunk in chunks:
            fp.write(json.dumps(chunk, ensure_ascii=True) + "\n")


def main() -> None:
    args = parse_args()
    raw_dir = Path(args.raw_dir)
    inventory_output = Path(args.inventory_output)
    chunks_output = Path(args.chunks_output)
    target_documents = {
        doc_id.strip().lower() for doc_id in args.target_documents.split(",") if doc_id.strip()
    }

    inventory = build_inventory(raw_dir)
    inventory, chunks, failed_documents = parse_targets(raw_dir, inventory, target_documents)

    write_inventory(inventory_output, inventory)
    write_chunks(chunks_output, chunks)

    parsed_docs = [entry["document_id"] for entry in inventory if entry.get("parse_status") == "parsed"]
    print(f"Inventory written: {inventory_output} ({len(inventory)} docs)")
    print(f"Chunks written: {chunks_output} ({len(chunks)} chunks)")
    print(f"Parsed documents: {parsed_docs}")
    if failed_documents:
        print("Failed documents:")
        for failed in failed_documents:
            print(f"- {failed['document_id']}: {failed['parse_status']} ({failed.get('parse_error', '')})")


if __name__ == "__main__":
    main()
